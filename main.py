# loading libraries
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
import seaborn as sns
import requests as rq
from pptx import Presentation
from PIL import Image
import io
import os
from bs4 import BeautifulSoup
import time

# loading data
vgchartz_data = pd.read_csv('datasets/pokemon_games.csv')

# printing initial data information
print(vgchartz_data.head())
print(vgchartz_data.info())
print(vgchartz_data.isnull().sum())

# sorting data by year
vgchartz_data_sorted = vgchartz_data.sort_values(by='Year', ascending=True)
print(vgchartz_data_sorted.tail())


# appending new data for an outdated dataset
new_data = {
    'Name': [
        'Pokémon Sword/Shield', 'Pokémon Scarlet/Violet', 'Pokémon Sun/Moon', 'Pokémon Ultra Sun/Ultra Moon',
        'Pokémon Yellow: Special Pikachu Edition', 'Pokémon Emerald Version', 'Pokémon Crystal Version',
        'Pokémon Platinum Version', 'Pokémon Legends: Arceus', 'New Pokémon Snap', 'Pokémon Rescue Team DX',
        'Pokémon Café ReMix', 'Pokémon Unite', 'Pokémon Trading Card Game', 'Pokémon Colosseum',
        'Pokémon Stadium 2', 'Pokémon Puzzle Challenge', 'Pokémon Omega Ruby/Alpha Sapphire',
        'Pokémon FireRed/LeafGreen', 'Pokémon Ranger series', 'Pokémon Let\'s Go, Pikachu!/Eevee!',
        'Pokémon Quest', 'Pokémon Masters EX', 'Pokémon Rumble Rush', 'Pokémon Puzzle League', 'Pokémon Channel',
        'Pokémon Legends: Celebi', 'Pokémon Mystery Dungeon: New Horizons', 'Pokémon Trading Card Game 2: Here Comes Team GR!'
    ],
    'Platform': [
        'NS', 'NS', '3DS', '3DS', 'GB', 'GBA', 'GBC', 'DS', 'NS', 'NS', 'NS', 'NS/Mobile', 'NS/Mobile', 'GB', 'GC',
        'N64', 'GBC', '3DS', 'GBA', 'DS', 'NS', 'NS/Mobile', 'Mobile', 'Mobile', 'N64', 'GC', 'NS', 'NS', 'GBC'
    ],
    'Year': [
        2019, 2022, 2016, 2017, 1998, 2004, 2000, 2008, 2022, 2021, 2020, 2020, 2021, 1998, 2003, 2000, 2000, 2014,
        2004, None, 2018, 2018, 2019, 2019, 2000, 2003, 2023, 2024, 2001
    ],
    'Genre': [
        'Role-Playing', 'Role-Playing', 'Role-Playing', 'Role-Playing', 'Role-Playing', 'Role-Playing', 'Role-Playing',
        'Role-Playing', 'Action/RPG', 'Simulation', 'Role-Playing', 'Puzzle', 'MOBA', 'Strategy', 'Role-Playing',
        'Strategy', 'Puzzle', 'Role-Playing', 'Role-Playing', 'Action-Adventure', 'Role-Playing', 'Action/Adventure',
        'Role-Playing', 'Action', 'Puzzle', 'Adventure', 'Action/RPG', 'Role-Playing', 'Strategy'
    ],
    'Publisher': [
        'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo',
        'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo',
        'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'DeNA', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo', 'Nintendo'
    ],
    'NA_Sales': [
        6.04, 5.92, 6.02, 2.46, 5.89, 2.57, 2.55, 2.76, 4.12, 1.26, 1.26, None, None, 1.49, 1.21, 1.02, None, 4.78,
        4.34, None, 4.17, None, None, None, 0.36, 0.24, None, None, 0.0
    ],
    'EU_Sales': [
        4.35, 4.12, 4.34, 1.74, 5.04, 1.58, 1.56, 1.73, 3.01, None, None, None, None, 0.73, 0.57, 0.36, None, 3.98,
        2.65, None, 3.02, None, None, None, 0.08, 0.06, None, None, 0.0
    ],
    'JP_Sales': [
        3.73, 3.12, 3.82, 2.55, 3.12, 2.06, 1.29, 2.69, 2.85, None, None, None, None, 1.38, 0.70, 1.13, None, 3.15,
        3.15, None, 2.98, None, None, None, None, 0.07, None, None, 0.89
    ],
    'Other_Sales': [
        1.12, 0.95, 0.99, 0.40, 0.59, 0.21, 0.99, 0.54, 0.95, None, None, None, None, 0.10, 0.07, 0.23, None, 0.82,
        0.35, None, 0.16, None, None, None, 0.01, 0.01, None, None, 0.0
    ],
    'Global_Sales': [
        15.24, 14.11, 15.17, 7.15, 14.64, 7.06, 6.39, 7.72, 10.93, 2.40, 1.26, None, None, 3.70, 2.54, 2.73, 0.45, 14.46,
        12.00, None, 14.33, None, None, None, 0.45, 0.38, None, None, 0.89
    ]
}
new_data_df = pd.DataFrame(new_data)

# combining and cleaning data
updated_data = pd.concat([vgchartz_data, new_data_df], ignore_index=True)
updated_data = updated_data.drop_duplicates(subset=['Name', 'Platform'], keep='last')

# saving updated data
updated_data.to_csv(file_path, index=False)
print("Dataset updated successfully!")

# printing missing values after update (optional)
print(vgchartz_data.isnull().sum())

# sorting data by year (ascending)
vgchartz_data_sorted = vgchartz_data.sort_values(by='Year', ascending=True)

# filtering data for finite global sales
vgchartz_data_sorted = vgchartz_data_sorted[np.isfinite(vgchartz_data_sorted['Global_Sales'])]

# configuring seaborn style
sns.set(style="whitegrid")

# creating a figure for the first bar plot
plt.figure(figsize=(14, 10))
barplot = sns.barplot(
    x='Name', 
    y='Global_Sales', 
    data=vgchartz_data_sorted, 
    palette='viridis'
)

# setting labels and title for the first bar plot
barplot.set_xlabel("Pokémon Game", fontsize=14)
barplot.set_ylabel("Global Sales (in millions)", fontsize=14)
barplot.set_title("Global Sales by Pokémon Games", fontsize=16)

# rotating x-axis labels for better readability and adjusting layout (first plot)
barplot.set_xticklabels(barplot.get_xticklabels(), rotation=90, ha='right', fontsize=12)
plt.tight_layout()
plt.show()

# moving onto selecting top 10 games based on global sales
top_10 = vgchartz_data_sorted.nlargest(10, "Global_Sales")

# configuring seaborn style (again)
sns.set(style="whitegrid")

# creating a figure for the second bar plot
plt.figure(figsize=(14, 8))

# generating bar plot of top 10 games
barplot = sns.barplot(
    x='Name', 
    y='Global_Sales', 
    data=top_10, 
    palette='rocket'
)
# setting labels and title for the second bar plot
barplot.set_xlabel("Game", fontsize=14)
barplot.set_ylabel("Global Sales (in millions)", fontsize=14)
barplot.set_title("Top 10 Pokémon Games by Global Sales", fontsize=16)

# rotating x-axis labels for better readability (second plot)
barplot.set_xticklabels(barplot.get_xticklabels(), rotation=45, ha='right', fontsize=12)

# adding data labels above bars in the second bar plot
for bar in barplot.patches:
    bar_height = bar.get_height()
    barplot.text(
        bar.get_x() + bar.get_width() / 2, 
        bar_height, 
        f'{bar_height:.2f}',  
        ha='center',  
        va='bottom',  
        fontsize=12,  
        color='black'
    )
plt.tight_layout()
plt.show()

# replacing missing sales values with 0
vgchartz_data[['NA_Sales', 'EU_Sales', 'JP_Sales', 'Other_Sales', 'Global_Sales']] = vgchartz_data[['NA_Sales', 'EU_Sales', 'JP_Sales', 'Other_Sales', 'Global_Sales']].fillna(0)

# calculating total sales by platform
platform_sales = vgchartz_data.groupby('Platform')['Global_Sales'].sum().reset_index()

# configuring seaborn style (2x again)
sns.set(style="whitegrid")

# creating a figure for the line plot
plt.figure(figsize=(16, 10))

# generating line plot of sales by platform
lineplot = sns.lineplot(
    data=platform_sales, 
    x='Platform', 
    y='Global_Sales', 
    marker='o'
)

# setting labels and title for the line plot (3rd plot)
lineplot.set_xlabel("Platform", fontsize=14)
lineplot.set_ylabel("Global Sales (in millions)", fontsize=14)
lineplot.set_title("Global Sales by Platform", fontsize=16)
plt.tight_layout()
plt.show()

# opening the presentation for TABLEAU data visualization
presentation = Presentation("tableau/best_viz.pptx")

# creating output directory for images
output_dir = "tableau_images"
os.makedirs(output_dir, exist_ok=True)

# extracting images from presentation slides
for slide_number, slide in enumerate(presentation.slides, start=1):
    for shape in slide.shapes:
        if shape.shape_type == 13:
            image = shape.image
            
            image_bytes = image.blob
            image_format = image.ext

            image_filename = os.path.join(output_dir, f"slide{slide_number}_image.{image_format}")
            with open(image_filename, "wb") as img_file:
                img_file.write(image_bytes)
            print(f"Saved image: {image_filename}")

# opening images for display
viz_1 = Image.open('tableau_images/slide1_image.png')
viz_2 = Image.open('tableau_images/slide2_image.png')

# displaying the first image (without axis)
plt.imshow(viz_1)
plt.axis('off')
plt.show()

# displaying the second image (without axis)
plt.imshow(viz_2)
plt.axis('off')
plt.show()

# web scraping attempt at vgchartz for pokemon game sales

# defining base URL and headers
BASE_URL = "https://www.vgchartz.com/games/games.php"
HEADERS = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}

# setting maximum number of pages to scrape
MAX_PAGES = 2

# defining search parameters for the website
params = {
    "name": "Pokemon",
    "order": "Sales",
    "ownership": "Both",
    "direction": "DESC",
    "showtotalsales": "1",  
    "shownasales": "1",    
    "showpalsales": "1",   
    "showjapansales": "1", 
    "showothersales": "1",  
    "showpublisher": "0",
    "showdeveloper": "0",
    "showreleasedate": "0",
    "showlastupdate": "0",
    "showvgchartzscore": "0",
    "showcriticscore": "0",
    "showuserscore": "0",
    "showshipped": "0"
}

# defining function to scrape data from a single page
def scrape_page(page_number):
    params["page"] = page_number
    print(f"Scraping page {page_number}...")
    response = requests.get(BASE_URL, headers=HEADERS, params=params)
    
    if response.status_code != 200:
        print(f"Failed to fetch page {page_number}. Status code: {response.status_code}")
        return []
    
    soup = BeautifulSoup(response.text, "html.parser")
    rows = soup.find_all("tr") 
    data = []

    for row in rows:
        cols = row.find_all("td")
        if len(cols) >= 8:  
            data.append({
                "Position": cols[0].get_text(strip=True),
                "Game": cols[1].get_text(strip=True),
                "Console": cols[2].get_text(strip=True),
                "Total Sales (M)": cols[3].get_text(strip=True),
                "NA Sales (M)": cols[4].get_text(strip=True),
                "EU Sales (M)": cols[5].get_text(strip=True),
                "JP Sales (M)": cols[6].get_text(strip=True),
                "Other Sales (M)": cols[7].get_text(strip=True),
            })
    return data

# defining function to scrape data from all specified pages
def scrape_all_pages():
    all_data = []
    for page in range(1, MAX_PAGES + 1):
        page_data = scrape_page(page)
        if not page_data:
            break 
        all_data.extend(page_data)
        time.sleep(1) 
    return all_data

# running the scraping process
if __name__ == "__main__":
    results = scrape_all_pages()
    if results:
        df = pd.DataFrame(results)
        df.to_csv("pokemon_game_sales.csv", index=False)
        print("Sales data scraped and saved to 'pokemon_game_sales.csv'.")
    else:
        print("No data scraped.")

# loading the scraped data along with the kaggle dataset from 8 years ago to merge
pokemon_sales_file = 'pokemon_game_sales.csv'
vgsales_file = 'vgsales.csv' #(https://www.kaggle.com/datasets/gregorut/videogamesales)

pokemon_sales_data = pd.read_csv(pokemon_sales_file)
vgsales_data = pd.read_csv(vgsales_file)

# worst case scenario: cleaning and formatting pokemon data
def clean_and_format_pokemon_data(pokemon_sales, vg_sales):
    rename_mapping = {
        "Position": "Rank",
        "Game": "Name",
        "Console": "Platform",
        "Total Sales (M)": "Global_Sales",
        "NA Sales (M)": "NA_Sales",
        "EU Sales (M)": "EU_Sales",
        "JP Sales (M)": "JP_Sales",
        "Other Sales (M)": "Other_Sales"
    }

    pokemon_sales = pokemon_sales.rename(columns=rename_mapping)

    common_columns = list(set(pokemon_sales.columns) & set(vg_sales.columns))

    pokemon_sales = pokemon_sales[common_columns].fillna(0)

    for column in ["Global_Sales", "NA_Sales", "EU_Sales", "JP_Sales", "Other_Sales"]:
        if column in pokemon_sales.columns:
            pokemon_sales[column] = (
                pokemon_sales[column]
                .replace(r'[^\d.]', '', regex=True)
                .replace('', '0')
                .astype(float)
            )

    return pokemon_sales

cleaned_pokemon_data = clean_and_format_pokemon_data(pokemon_sales_data, vgsales_data)

# saving cleaned data to a new CSV file
output_path = 'pokemon_game_sales_cleaned_v2.csv'
cleaned_pokemon_data.to_csv(output_path, index=False)

# API interaction for pokemon data

# displaying pokemon images for top 2 pokemon of all time (opinions vary)
switch = "resources/ns.jpg" #https://collider.com/most-popular-pokemon-characters-ranked/
gameboy = "resources/gb.jpg" #https://bulbapedia.bulbagarden.net/wiki/Pok%C3%A9mon_of_the_Year

plt.figure(figsize=(5, 5))
electric = mpimg.imread(switch)
plt.imshow(electric)
plt.axis("off")
plt.title("Pikachu: The Face of Pokémon")
plt.show()

plt.figure(figsize=(5, 5))
fire = mpimg.imread(gameboy)
plt.imshow(fire)
plt.axis("off")
plt.title("Charizard: The Powerful Flame")
plt.show()

# delving deeper into the pokemon game world by talking about the top two most popular pokemon

# OOP approach to fetching pokemon data
class PokeAPI:
    # base url for accessing pokemon data on pokeapi.co
    BASE_URL = "https://pokeapi.co/api/v2/"
    
    def __init__(self):
        pass

    def get_pokemon_data(self, name_or_id):
        # constructing the url for fetching pokemon data
        url = f"{self.BASE_URL}pokemon/{name_or_id}"
        try:
            # sending a get request to the pokemon api url
            response = requests.get(url)
            # raising an exception if the request fails
            response.raise_for_status()
            # returning the json response as a dictionary
            return response.json()
        except requests.exceptions.RequestException as e:
            # printing an error message if data fetching fails
            print(f"Error fetching data: {e}")
            # returning none if data fetching fails
            return None

    def extract_pokemon_stats(self, pokemon_data):
        # checking if pokemon data exists
        if pokemon_data:
            # extracting the stats as a dictionary
            stats = {stat['stat']['name']: stat['base_stat'] for stat in pokemon_data['stats']}
            # returning the extracted pokemon stats dictionary
            return stats
        # returning none if pokemon data does not exist
        return None

# visualizing pokemon stats    
def visualize_stats(pokemon_name, stats, title_color, bar_color):
    if stats:
        plt.figure(figsize=(8, 5))
        sns.barplot(x=list(stats.keys()), y=list(stats.values()), palette=[bar_color])
        plt.title(f"{pokemon_name.capitalize()}'s Stats", fontsize=16, color=title_color)
        plt.xlabel("Stat", fontsize=12)
        plt.ylabel("Base Stat Value", fontsize=12)
        plt.xticks(rotation=45, fontsize=10)
        plt.grid(axis="y", linestyle="--", alpha=0.6)
        plt.show()

# creating an instance of the PokeAPI class
api = PokeAPI()

# fetching data for pikachu
pikachu_data = api.get_pokemon_data("pikachu")
# extracting stats for pikachu
pikachu_stats = api.extract_pokemon_stats(pikachu_data)

# visualizing pikachu's stats
visualize_stats("Pikachu", pikachu_stats, title_color="black", bar_color="gold")

# fetching data for charizard
charizard_data = api.get_pokemon_data("charizard")
# extracting stats for charizard
charizard_stats = api.extract_pokemon_stats(charizard_data)

# visualizing charizard's stats
visualize_stats("Charizard", charizard_stats, title_color="orangered", bar_color="orangered")

# comparing pokemon stats
def compare_pokemon_stats(pikachu_stats, charizard_stats):
    df = pd.DataFrame({
        "Stat": list(pikachu_stats.keys()),
        "Pikachu": list(pikachu_stats.values()),
        "Charizard": list(charizard_stats.values())
    })

    # melting the dataframe for seaborn plotting
    df = df.melt(id_vars="Stat", var_name="Pokemon", value_name="Base Stat")
    plt.figure(figsize=(10, 6))
    sns.barplot(data=df, x="Stat", y="Base Stat", hue="Pokemon", palette=["gold", "orangered"])
    plt.title("Comparison of Pikachu and Charizard Stats", fontsize=16)
    plt.xlabel("Stat", fontsize=12)
    plt.ylabel("Base Stat Value", fontsize=12)
    plt.xticks(rotation=45, fontsize=10)
    plt.legend(title="Pokemon")
    plt.grid(axis="y", linestyle="--", alpha=0.6)
    plt.show()

# comparing pikachu and charizard stats
compare_pokemon_stats(pikachu_stats, charizard_stats)

# a fun program that generates stats for YOUR favorite pokemon!

# creating a dictionary of pokemon types and their corresponding colors
pokemon_type_colors = {
    "normal": "#A8A77A",  # Grayish
    "fire": "#EE8130",    # Orange-red
    "water": "#6390F0",   # Blue
    "grass": "#7AC74C",   # Green
    "electric": "#F7D02C", # Yellow
    "psychic": "#F95587",  # Pink
    "ice": "#96D9D6",     # Light blue
    "fighting": "#C22E28", # Reddish-brown
    "poison": "#A33EA1",  # Purple
    "ground": "#E2BF65",  # Brown
    "flying": "#A98FF3",  # Light purple
    "bug": "#A6B91A",     # Greenish-yellow
    "rock": "#B6A136",    # Dark yellow
    "ghost": "#735797",   # Dark purple
    "dragon": "#6F35FC",  # Purple-blue
    "steel": "#B7B7CE",   # Light gray
    "dark": "#705746",    # Dark brown
    "fairy": "#D685AD",   # Pinkish-purple
}

# getting pokemon name from user input
def get_pokemon_name():
    # entering pokemon name
    while True:
        pokemon_name = input("Enter a Pokemon name (or 'q' to quit): ").lower()
        # checking for quit condition
        if pokemon_name == 'q':
            break
        return pokemon_name

# looping for multiple pokemon visualization
while True:
    pokemon_name = get_pokemon_name()
    # checking if pokemon name is valid
    if not pokemon_name:
        break

    # fetching pokemon data
    pokemon_data = api.get_pokemon_data(pokemon_name)

    # checking if data was retrieved successfully
    if pokemon_data:
        # extracting pokemon stats
        pokemon_stats = api.extract_pokemon_stats(pokemon_data)
        types = [t['type']['name'] for t in pokemon_data['types']]

        # choosing color based on pokemon types
        if types:
            bar_color = pokemon_type_colors.get(types[0], "gray")  # Use first type or gray if not found
            if len(types) > 1:
                title_color = pokemon_type_colors.get(types[1],"black")
            else:
                title_color = pokemon_type_colors.get(types[0],"black")
        else:
            bar_color = "gray"
            title_color = "black"

        # visualizing pokemon's stats
        visualize_stats(pokemon_name.capitalize(), pokemon_stats, title_color, bar_color)
    else:
        # printing error message if data not found
        print(f"Sorry, couldn't find data for '{pokemon_name}'.")

# printing completion message
print("Done visualizing Pokemon stats!")

# ~end of main.py~